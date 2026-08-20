# -*- coding: utf-8 -*-
"""生成《全球公共卫生监测与口岸预警系统》双平台介绍 PPT。"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ---------- 设计令牌（暖纸学术风，与项目 UI 一致） ----------
INK = RGBColor(0x2A, 0x33, 0x3C)        # 深墨
CRIMSON = RGBColor(0x9D, 0x30, 0x39)    # 绛红主色
TEAL = RGBColor(0x1F, 0x7A, 0x6E)       # 青绿强调
PAPER = RGBColor(0xF6, 0xF4, 0xEE)      # 暖纸底
CARD = RGBColor(0xFF, 0xFF, 0xFF)
MIST = RGBColor(0xE9, 0xE5, 0xDB)       # 卡片描边灰
GRAY = RGBColor(0x6B, 0x77, 0x82)       # 正文灰
LIGHT = RGBColor(0xF3, 0xF0, 0xE8)
NAVY = RGBColor(0x2F, 0x44, 0x54)
GOLD = RGBColor(0xB8, 0x8A, 0x2E)

FONT_TITLE = "Noto Serif SC"
FONT_BODY = "Noto Sans SC"
FONT_LATIN = "Georgia"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def set_font(run, size, color=INK, bold=False, font=FONT_BODY, latin=None, spacing=None):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = latin or font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)
    if spacing is not None:
        rPr.set('spc', str(spacing))

def rect(slide, x, y, w, h, fill=CARD, line=None, radius=None, shadow=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radius is not None:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape

def textbox(slide, x, y, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """lines: list of (text, size, color, bold, font, space_after)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, item in enumerate(lines):
        text, size, color, bold, font, after = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if after:
            p.space_after = Pt(after)
        run = p.add_run()
        run.text = text
        set_font(run, size, color, bold, font)
    return tb

def slide_base(fill=PAPER):
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    return slide

def header(slide, eyebrow, title, sub=None):
    rect(slide, Inches(0.55), Inches(0.42), Inches(0.09), Inches(0.86), fill=CRIMSON)
    textbox(slide, Inches(0.82), Inches(0.38), Inches(11.6), Inches(0.3),
            [(eyebrow, 10.5, CRIMSON, True, FONT_BODY, 0)])
    textbox(slide, Inches(0.82), Inches(0.62), Inches(11.8), Inches(0.55),
            [(title, 27, INK, True, FONT_TITLE, 0)])
    if sub:
        textbox(slide, Inches(0.82), Inches(1.22), Inches(11.8), Inches(0.3),
                [(sub, 11.5, GRAY, False, FONT_BODY, 0)])
    return slide

def page_no(slide, n, total):
    textbox(slide, SW - Inches(1.15), SH - Inches(0.42), Inches(0.75), Inches(0.25),
            [(f"{n:02d} / {total}", 9, GRAY, False, FONT_LATIN, 0)], align=PP_ALIGN.RIGHT)

def chip(slide, x, y, w, h, text, fill=LIGHT, fg=NAVY, size=10.5, bold=True):
    rect(slide, x, y, w, h, fill=fill, radius=0.5)
    textbox(slide, x, y + Emu(int(h/2) - Pt(size*1.05)), w, Pt(size*1.4),
            [(text, size, fg, bold, FONT_BODY, 0)], align=PP_ALIGN.CENTER)

def card(slide, x, y, w, h, title, body_lines, icon_no=None, accent=CRIMSON, title_size=13.5, body_size=10.5):
    rect(slide, x, y, w, h, fill=CARD, line=MIST, radius=0.07)
    rect(slide, x, y, Inches(0.055), h, fill=accent)
    ty = y + Inches(0.16)
    if icon_no:
        textbox(slide, x + Inches(0.22), ty, Inches(0.5), Inches(0.5),
                [(icon_no, 20, RGBColor(0xD8, 0xD2, 0xC6), True, FONT_LATIN, 0)])
        tx = x + Inches(0.78)
    else:
        tx = x + Inches(0.24)
    textbox(slide, tx, ty, w - Inches(0.95), Inches(0.35),
            [(title, title_size, INK, True, FONT_TITLE, 0)])
    lines = [(t, body_size, GRAY, False, FONT_BODY, 3.5) for t in body_lines]
    textbox(slide, tx, ty + Inches(0.42), w - Inches(1.0), h - Inches(0.6), lines)

TOTAL = 14
N = 0
def next_no():
    global N
    N += 1
    return N, TOTAL

# ============================================================ 封面
s = slide_base(PAPER)
rect(s, 0, 0, SW, SH, fill=PAPER)
rect(s, Inches(9.4), 0, Inches(3.933), SH, fill=NAVY)
rect(s, 0, Inches(7.28), SW, Inches(0.22), fill=CRIMSON)
# 左侧主标题区
textbox(s, Inches(0.9), Inches(1.5), Inches(8.2), Inches(0.35),
        [("INTERNATIONAL PUBLIC HEALTH SURVEILLANCE & PORT ALERT SYSTEM", 11, CRIMSON, True, FONT_LATIN, 0)], )
textbox(s, Inches(0.9), Inches(1.95), Inches(8.3), Inches(1.9),
        [("全球公共卫生监测", 44, INK, True, FONT_TITLE, 6),
         ("与口岸预警系统", 44, INK, True, FONT_TITLE, 0)])
rect(s, Inches(0.92), Inches(3.95), Inches(1.5), Inches(0.045), fill=CRIMSON)
textbox(s, Inches(0.9), Inches(4.25), Inches(8.2), Inches(1.2),
        [("互联网监测平台 · 内网口岸预警平台 —— 一体化双平台架构", 15, NAVY, True, FONT_BODY, 8),
         ("公开数据采集、风险研判、单向摆渡，到旅客风险匹配与口岸布控的全链路闭环", 12, GRAY, False, FONT_BODY, 0)])
chip(s, Inches(0.9), Inches(5.5), Inches(1.75), Inches(0.42), "数据采集", LIGHT, NAVY)
chip(s, Inches(2.8), Inches(5.5), Inches(1.75), Inches(0.42), "风险研判", LIGHT, NAVY)
chip(s, Inches(4.7), Inches(5.5), Inches(1.75), Inches(0.42), "单向摆渡", LIGHT, NAVY)
chip(s, Inches(6.6), Inches(5.5), Inches(1.95), Inches(0.42), "口岸布控建议", LIGHT, NAVY)
textbox(s, Inches(0.9), Inches(6.55), Inches(8.2), Inches(0.3),
        [("InternationalPublicHealth · 项目全景介绍", 10.5, GRAY, False, FONT_BODY, 0)])
# 右侧深蓝区装饰
textbox(s, Inches(9.9), Inches(1.6), Inches(3.1), Inches(2.2),
        [("两朵云", 22, PAPER, True, FONT_TITLE, 4),
         ("一套理念", 15, RGBColor(0x9F, 0xB4, 0xC2), False, FONT_BODY, 10),
         ("互联网采 全球看", 13, RGBColor(0x9F, 0xB4, 0xC2), False, FONT_BODY, 4),
         ("内网测 口岸拦", 13, RGBColor(0x9F, 0xB4, 0xC2), False, FONT_BODY, 0)])
for i, (yy, ww, cc) in enumerate([(Inches(4.3), Inches(2.6), CRIMSON), (Inches(4.85), Inches(2.0), TEAL), (Inches(5.4), Inches(1.5), GOLD)]):
    rect(s, Inches(9.9), yy, ww, Inches(0.05), fill=cc)
n, _ = next_no()
textbox(s, Inches(9.9), Inches(6.55), Inches(2.9), Inches(0.3),
        [("2026 · 项目组", 10.5, RGBColor(0x9F, 0xB4, 0xC2), False, FONT_BODY, 0)])

# ============================================================ 2 总览
s = slide_base()
header(s, "OVERVIEW", "系统总览：一体化双平台，一张蓝图", "依据《001-需求》建设两个逻辑与物理完全隔离的平台，覆盖“全球监测 → 风险研判 → 单向摆渡 → 口岸拦截”全链路")
# 中间数据流
flow = [("公开数据源", NAVY), ("互联网采集·清洗", CRIMSON), ("风险/规则/地图", CRIMSON), ("加密签名单向数据包", GOLD), ("内网接收·验签·解密", TEAL), ("旅客匹配·预警·布控", TEAL)]
fx = Inches(0.75); fw = Inches(1.85); gap = Inches(0.13)
for i, (t, c) in enumerate(flow):
    x = fx + i * (fw + gap)
    rect(s, x, Inches(1.85), fw, Inches(0.75), fill=c, radius=0.12)
    textbox(s, x + Inches(0.06), Inches(2.03), fw - Inches(0.12), Inches(0.5),
            [(t, 11.5, PAPER, True, FONT_BODY, 0)], align=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        textbox(s, x + fw - Inches(0.02), Inches(2.08), gap + Inches(0.05), Inches(0.3),
                [("▸", 14, GRAY, True, FONT_LATIN, 0)], align=PP_ALIGN.CENTER)
# 三大理念卡
card(s, Inches(0.75), Inches(3.05), Inches(3.85), Inches(2.5),
     "面向全球公共卫生治理",
     ["对 WHO、ECDC、ProMED-mail、JHU CSSE、OWID、HealthMap 六类权威公开源做自动化采集与结构化",
      "以六因子风险引擎量化各国疫情风险，形成四级阈值与预警",
      "以可视化地图支撑全球态势研判与决策"],
     icon_no="Ⅰ", accent=CRIMSON)
card(s, Inches(4.75), Inches(3.05), Inches(3.85), Inches(2.5),
     "面向口岸一线实战",
     ["内网落地旅客风险匹配：14 天旅居史、健康申报、中转链路实时比对",
      "按机场/海港/陆路生成结构化布控措施与法律依据",
      "支持单人录入、批量导入、目录自动处理与数据库直连同步"],
     icon_no="Ⅱ", accent=TEAL)
card(s, Inches(8.75), Inches(3.05), Inches(3.85), Inches(2.5),
     "安全即架构",
     ["互联网与内网不共享数据库、Redis、JWT 密钥与字段密钥",
      "JWT audience 双向隔离，API 路由按平台裁剪",
      "摆渡包认证加密 + 签名 + 幂等，等保三级技术基线"],
     icon_no="Ⅲ", accent=GOLD)
# 底部数字条
rect(s, Inches(0.75), Inches(5.8), Inches(11.85), Inches(0.95), fill=CARD, line=MIST, radius=0.1)
stats = [("6", "权威公开数据源"), ("6", "风险研判因子"), ("4", "风险等级阈值"), ("3", "摆渡通道"), ("5", "RBAC 角色"), ("2", "独立隔离平台")]
sx = Inches(1.05)
for v, label in stats:
    textbox(s, sx, Inches(5.95), Inches(1.8), Inches(0.45),
            [(v, 24, CRIMSON, True, FONT_LATIN, 0)], align=PP_ALIGN.CENTER)
    textbox(s, sx, Inches(6.42), Inches(1.8), Inches(0.25),
            [(label, 9.5, GRAY, False, FONT_BODY, 0)], align=PP_ALIGN.CENTER)
    sx += Inches(1.92)
page_no(s, *next_no())

# ============================================================ 3 互联网定位
s = slide_base()
header(s, "INTERNET PLATFORM", "互联网监测平台：定位与目标", "外网威胁感知的前哨站 —— 持续扫描全球公开疫情信号，转化为可研判、可交付的风险情报")
card(s, Inches(0.75), Inches(1.75), Inches(5.8), Inches(2.15), "平台定位",
     ["部署于互联网安全域，面向数据分析员与只读用户",
      "“采集 → 清洗 → 评分 → 研判 → 摆渡”的情报生产线",
      "不接触任何旅客个人信息，专注公开数据"],
     icon_no="01", accent=CRIMSON)
card(s, Inches(6.8), Inches(1.75), Inches(5.8), Inches(2.15), "建设目标",
     ["全球疫情态势分钟级感知，异常趋势自动预警",
      "风险结论标准化、可解释、可追溯到原始来源",
      "把“看新闻”升级为“看仪表盘”"],
     icon_no="02", accent=TEAL)
# 六大功能模块
mods = [
    ("数据采集", "六源适配器 · 调度热更新 · 退避重试 · 指纹去重", CRIMSON),
    ("疫情事件", "等级过滤 · 关键词检索 · CSV 导出 · 模型研判", CRIMSON),
    ("风险研判", "六因子权重 · 国家排行 · 区域趋势 · 预警管理", CRIMSON),
    ("全球展示", "MapLibre+Deck.gl · 六类图层 · 14 日时间轴 · 下钻导出", CRIMSON),
    ("数据摆渡", "文件/MQ/outbox 三通道 · 加密签名 · 分片续传", GOLD),
    ("规则引擎", "条件树 all/any/not · 版本草稿发布 · 在线测试 · 热更新", CRIMSON),
]
mx = Inches(0.75); my = Inches(4.15); mw = Inches(3.8); mh = Inches(1.28)
for i, (t, d, c) in enumerate(mods):
    x = mx + (i % 3) * (mw + Inches(0.15))
    y = my + (i // 3) * (mh + Inches(0.15))
    rect(s, x, y, mw, mh, fill=CARD, line=MIST, radius=0.09)
    rect(s, x, y, Inches(0.05), mh, fill=c)
    textbox(s, x + Inches(0.2), y + Inches(0.13), mw - Inches(0.35), Inches(0.3),
            [(t, 13, INK, True, FONT_TITLE, 0)])
    textbox(s, x + Inches(0.2), y + Inches(0.5), mw - Inches(0.35), Inches(0.7),
            [(d, 9.5, GRAY, False, FONT_BODY, 0)])
page_no(s, *next_no())

# ============================================================ 4 互联网功能细节
s = slide_base()
header(s, "INTERNET PLATFORM", "互联网端核心能力详解", "从原始信号到风险情报的完整加工链")
# 左列：采集与解析
card(s, Inches(0.75), Inches(1.7), Inches(3.85), Inches(2.6), "六源采集与用户配置",
     ["信息源全部数据库化：URL、启停、固定间隔 / 五段 Cron、IANA 时区",
      "解析方式三选一：内置结构化 / 大模型 / 混合",
      "预设源仅首次空库导入，运行期完全由用户掌控",
      "源失败独立记录，不阻断其他源"],
     icon_no="A", accent=CRIMSON)
# 中列：风险引擎
card(s, Inches(4.75), Inches(1.7), Inches(3.85), Inches(2.6), "六因子风险引擎",
     ["权重合计强制 100%，保存即生成新草稿",
      "国家历史、趋势突变参与计算",
      "四级阈值：低 / 中 / 高 / 极高",
      "草稿 → 规则页测试 → 发布 → 即时热更新"],
     icon_no="B", accent=TEAL)
# 右列：地图
card(s, Inches(8.75), Inches(1.7), Inches(3.85), Inches(2.6), "全球态势地图",
     ["MapLibre GL + Deck.gl，本地边界数据无 CDN",
      "六类图层：风险 / 气泡 / 热力 / 中转 / 人员流 / 内网口岸",
      "真实 14 日时间轴拖动回放",
      "点击国家下钻，导出 PNG 报告"],
     icon_no="C", accent=GOLD)
# 底部：摆渡三通道
rect(s, Inches(0.75), Inches(4.55), Inches(11.85), Inches(2.25), fill=CARD, line=MIST, radius=0.06)
textbox(s, Inches(1.05), Inches(4.78), Inches(4.0), Inches(0.35),
        [("单向数据摆渡 · 三通道", 14, INK, True, FONT_TITLE, 0)])
ch_items = [
    ("文件通道", "网闸目录投递，gzip + SHA-256 完整性校验"),
    ("消息队列", "RabbitMQ 受控投递，分片续传与幂等版本入库"),
    ("API Outbox", "HTTPS 白名单网关，全量 / 增量任务"),
]
cx = Inches(1.05)
for t, d in ch_items:
    rect(s, cx, Inches(5.3), Inches(3.7), Inches(1.2), fill=LIGHT, radius=0.1)
    textbox(s, cx + Inches(0.2), Inches(5.45), Inches(3.3), Inches(0.3),
            [(t, 12, NAVY, True, FONT_BODY, 0)])
    textbox(s, cx + Inches(0.2), Inches(5.8), Inches(3.3), Inches(0.6),
            [(d, 9.5, GRAY, False, FONT_BODY, 0)])
    cx += Inches(3.9)
textbox(s, Inches(1.05), Inches(6.55), Inches(11.2), Inches(0.3),
        [("密码套件：AES-256-GCM + Ed25519 签名，或 SM4-CBC + SM2-SM3 国密双套件可选", 10, CRIMSON, True, FONT_BODY, 0)])
page_no(s, *next_no())

# ============================================================ 5 内网定位
s = slide_base()
header(s, "INTRANET PLATFORM", "内网口岸预警平台：定位与目标", "口岸一线的拦截防线 —— 在物理隔离环境中完成旅客风险匹配与布控决策")
card(s, Inches(0.75), Inches(1.75), Inches(5.8), Inches(2.15), "平台定位",
     ["部署于内网安全域（Docker internal 网络），面向口岸操作员",
      "离线镜像展示 + 摆渡接收 + 旅客研判的作战平台",
      "与互联网展示同源同构，另增口岸图层与旅客预警"],
     icon_no="01", accent=TEAL)
card(s, Inches(6.8), Inches(1.75), Inches(5.8), Inches(2.15), "建设目标",
     ["旅客落地即研判：14 天旅居史 × 国家风险 × 规则引擎实时匹配",
      "布控建议附带结构化措施与法律依据，可直接执行",
      "敏感字段全程加密脱敏，明文仅在处理内存中出现"],
     icon_no="02", accent=CRIMSON)
mods = [
    ("旅客风险预警", "单人录入 · CSV/JSONL 批量导入 · 目录自动处理 · 数据库直连同步", TEAL),
    ("布控建议", "按机场 / 海港 / 陆路区分的结构化措施 + 法律依据", TEAL),
    ("数据接收中心", "目录扫描 · RabbitMQ · 人工包导入 · API 轮询（默认关闭）", TEAL),
    ("离线镜像地图", "与互联网一致的全球态势，本地渲染无外联", TEAL),
    ("规则与风险", "内网独立的规则版本管理与风险重算能力", TEAL),
    ("认证与管理", "LDAP/AD + 本地账号 · 管理员强制 MFA · 审计留痕", TEAL),
]
mx = Inches(0.75); my = Inches(4.15); mw = Inches(3.8); mh = Inches(1.28)
for i, (t, d, c) in enumerate(mods):
    x = mx + (i % 3) * (mw + Inches(0.15))
    y = my + (i // 3) * (mh + Inches(0.15))
    rect(s, x, y, mw, mh, fill=CARD, line=MIST, radius=0.09)
    rect(s, x, y, Inches(0.05), mh, fill=c)
    textbox(s, x + Inches(0.2), y + Inches(0.13), mw - Inches(0.35), Inches(0.3),
            [(t, 13, INK, True, FONT_TITLE, 0)])
    textbox(s, x + Inches(0.2), y + Inches(0.5), mw - Inches(0.35), Inches(0.7),
            [(d, 9.5, GRAY, False, FONT_BODY, 0)])
page_no(s, *next_no())

# ============================================================ 6 旅客数据链路
s = slide_base()
header(s, "PASSENGER PIPELINE", "内网端：旅客风险数据链路", "多渠道汇聚 → 字段映射 → 风险匹配 → 布控建议的闭环")
# 顶部：四渠道
chans = [
    ("单人录入", "证件、姓名、国籍、14 日旅居史、中转、口岸、航班、健康申报", TEAL),
    ("批量导入", "UTF-8 CSV / JSONL 文件，导入即返回脱敏结果", TEAL),
    ("目录自动处理", "口岸系统投递规范文件至配置目录，后台自动解析入库", TEAL),
    ("数据库直连同步", "MySQL / PostgreSQL / SQL Server / Oracle / ClickHouse / 阿里云 MaxCompute(ODPS) / AnalyticDB(ADS)", CRIMSON),
]
cx = Inches(0.75)
for t, d, c in chans:
    rect(s, cx, Inches(1.75), Inches(2.86), Inches(1.55), fill=CARD, line=MIST, radius=0.09)
    rect(s, cx, Inches(1.75), Inches(2.86), Inches(0.07), fill=c)
    textbox(s, cx + Inches(0.18), Inches(1.95), Inches(2.5), Inches(0.3),
            [(t, 12.5, INK, True, FONT_TITLE, 0)])
    textbox(s, cx + Inches(0.18), Inches(2.32), Inches(2.55), Inches(0.9),
            [(d, 9.5, GRAY, False, FONT_BODY, 0)])
    cx += Inches(3.0)
# 中部流程
steps = [("智能字段映射", "中英文列名自动识别，模糊匹配兜底"), ("旅客入库", "复用统一录入链路，自动触发风险匹配"), ("风险引擎匹配", "旅居史 × 国家风险 × 规则树实时计算"), ("布控建议生成", "分交通方式的措施与法律依据")]
fx = Inches(0.75); fw = Inches(2.72)
for i, (t, d) in enumerate(steps):
    x = fx + i * (fw + Inches(0.32))
    rect(s, x, Inches(3.7), fw, Inches(1.0), fill=NAVY, radius=0.12)
    textbox(s, x + Inches(0.12), Inches(3.83), fw - Inches(0.24), Inches(0.3),
            [(t, 12, PAPER, True, FONT_BODY, 0)], align=PP_ALIGN.CENTER)
    textbox(s, x + Inches(0.12), Inches(4.18), fw - Inches(0.24), Inches(0.5),
            [(d, 8.5, RGBColor(0x9F, 0xB4, 0xC2), False, FONT_BODY, 0)], align=PP_ALIGN.CENTER)
    if i < 3:
        textbox(s, x + fw, Inches(3.95), Inches(0.32), Inches(0.35),
                [("▸", 16, GRAY, True, FONT_LATIN, 0)], align=PP_ALIGN.CENTER)
# 底部：数据库连接能力
rect(s, Inches(0.75), Inches(5.0), Inches(11.85), Inches(1.75), fill=CARD, line=MIST, radius=0.06)
textbox(s, Inches(1.05), Inches(5.2), Inches(6.0), Inches(0.35),
        [("广泛数据库连接能力（对齐 DeepAnalyze）", 13.5, INK, True, FONT_TITLE, 0)])
db_rows = [
    ("阿里云 MaxCompute / ODPS", "pyodps 专用驱动 · Endpoint + Project + AccessKey · Tunnel 直读", CRIMSON),
    ("阿里云 AnalyticDB (ADS)", "MySQL 兼容 / PostgreSQL 兼容双协议接入", CRIMSON),
    ("关系型与分析型数据库", "MySQL · RDS · PolarDB · GBase · PostgreSQL · SQL Server · Oracle · ClickHouse", TEAL),
]
dy = Inches(5.62)
for t, d, c in db_rows:
    rect(s, Inches(1.05), dy + Inches(0.07), Inches(0.09), Inches(0.28), fill=c)
    textbox(s, Inches(1.3), dy, Inches(3.9), Inches(0.3), [(t, 10.5, NAVY, True, FONT_BODY, 0)])
    textbox(s, Inches(5.3), dy, Inches(7.1), Inches(0.3), [(d, 9.5, GRAY, False, FONT_BODY, 0)])
    dy += Inches(0.38)
page_no(s, *next_no())

# ============================================================ 7 大模型能力
s = slide_base()
header(s, "LLM CAPABILITY", "大语言模型能力：双侧独立，安全可控", "统一网关适配五类协议，配置独立落库、密钥独立加密，不经摆渡通道传输")
# 协议条
protos = ["OpenAI", "OpenAI 兼容", "Anthropic", "Google Gemini", "Ollama 本地"]
px = Inches(0.75)
for t in protos:
    rect(s, px, Inches(1.75), Inches(2.25), Inches(0.5), fill=NAVY, radius=0.3)
    textbox(s, px, Inches(1.86), Inches(2.25), Inches(0.3),
            [(t, 11, PAPER, True, FONT_BODY, 0)], align=PP_ALIGN.CENTER)
    px += Inches(2.4)
textbox(s, Inches(0.75), Inches(2.35), Inches(11.9), Inches(0.28),
        [("丰富提供商预设：OpenAI · Claude · Gemini · DeepSeek · 阿里百炼/Qwen · 智谱 GLM · Kimi · 豆包/火山方舟 · 百度千帆 · SiliconFlow · Ollama · 自定义兼容", 10, CRIMSON, True, FONT_BODY, 0)])
# 左：管理能力
card(s, Inches(0.75), Inches(2.85), Inches(5.8), Inches(2.35), "全生命周期管理",
     ["获取模型：拉取供应商模型列表，一键选择",
      "连续测试：逐模型测试连通性并显示连接时长（ms）",
      "单模型测试：真实最小对话验证",
      "默认模型：每侧独立指定，互不影响",
      "密钥加密落库，API 只回传 has_api_key"],
     icon_no="⚙", accent=CRIMSON)
# 右：应用场景
card(s, Inches(6.8), Inches(2.85), Inches(5.8), Inches(2.35), "两大应用场景",
     ["事件研判：疫情事件详情页调用本平台默认模型辅助分析，区分事实与建议",
      "信息提取：采集源选用大模型/混合解析，自定义提取提示词",
      "输出仅作辅助研判，不替代来源事实",
      "保留原始采集文件以供核验"],
     icon_no="◎", accent=TEAL)
# 底部安全边界
rect(s, Inches(0.75), Inches(5.45), Inches(11.85), Inches(1.3), fill=LIGHT, radius=0.08)
textbox(s, Inches(1.05), Inches(5.65), Inches(11.3), Inches(0.9),
        [("安全边界", 12.5, NAVY, True, FONT_TITLE, 5),
         ("两侧配置分别存于各自数据库，用各自 FIELD_ENCRYPTION_KEY 加密；供应商地址仅接受无账号无参数的 HTTP(S)，禁云元数据地址，不跟随重定向；内网默认零配置零外联，物理隔离场景指向内网 Ollama / 推理服务。", 10.5, GRAY, False, FONT_BODY, 0)])
page_no(s, *next_no())

# ============================================================ 8 安全设计
s = slide_base()
header(s, "SECURITY", "安全设计：等保三级技术基线", "安全不是附加项，而是架构本身")
sec_cards = [
    ("身份鉴别", ["bcrypt 独立盐 + 密码复杂度", "5 次失败锁定 15 分钟", "TOTP 双因素，内网管理员强制 MFA", "LDAP/AD 与本地故障账号并存"], CRIMSON),
    ("访问控制", ["五角色最小权限 RBAC", "审计员与系统管理员职责分离", "互联网/内网 JWT audience 双向隔离", "路由与构建产物按平台裁剪"], TEAL),
    ("数据保护", ["旅客证件、姓名、联系方式 AES-256-GCM 字段加密", "证件盲索引 + 页面脱敏", "日志不记录请求体", "明文仅存在于处理内存"], GOLD),
    ("边界与摆渡", ["内网 Docker 网络 internal，前端 CSP 严格", "摆渡包认证加密 + 独立签名 + SHA-256", "Schema 校验 + package ID 幂等", "API 轮询默认关闭"], NAVY),
    ("审计与可用性", ["全部写操作数据库留痕", "JSONL 应用访问日志 + 请求 ID", "PostgreSQL WAL + 每日备份 14 日保留", "Redis 故障降级进程缓存"], CRIMSON),
]
cx = Inches(0.75); cy = Inches(1.8)
cw = Inches(3.85); chh = Inches(2.35)
for i, (t, items, c) in enumerate(sec_cards):
    x = cx + (i % 3) * (cw + Inches(0.15))
    y = cy + (i // 3) * (chh + Inches(0.18))
    if i == 4:
        x = cx; y = cy + chh + Inches(0.18); 
    rect(s, x, y, cw, chh, fill=CARD, line=MIST, radius=0.07)
    rect(s, x, y, cw, Inches(0.06), fill=c)
    textbox(s, x + Inches(0.22), y + Inches(0.18), cw - Inches(0.4), Inches(0.32),
            [(t, 13, INK, True, FONT_TITLE, 0)])
    lines = [("· " + t2, 9.5, GRAY, False, FONT_BODY, 3) for t2 in items]
    textbox(s, x + Inches(0.22), y + Inches(0.58), cw - Inches(0.42), chh - Inches(0.7), lines)
# 第六格放国密说明
x = cx + cw + Inches(0.15); y = cy + chh + Inches(0.18)
rect(s, x, y, cw * 2 + Inches(0.15), chh, fill=NAVY, radius=0.07)
textbox(s, x + Inches(0.25), y + Inches(0.18), Inches(7.2), Inches(0.32),
        [("国密双套件与密钥职责分离", 13, PAPER, True, FONT_TITLE, 0)])
textbox(s, x + Inches(0.25), y + Inches(0.6), Inches(7.3), Inches(1.6),
        [("SM2 加解密与签名各持一套密钥对：互联网只持接收方公钥用于包封，内网只持私钥解封；签名私钥仅在互联网，验签公钥在内网。", 10, RGBColor(0x9F, 0xB4, 0xC2), False, FONT_BODY, 5),
         ("生产密钥从部署密钥注入，不写入代码、日志与备份资源清单；可用 AES-256-GCM + Ed25519 或 SM4-CBC + SM2-SM3 双选。", 10, RGBColor(0x9F, 0xB4, 0xC2), False, FONT_BODY, 0)])
page_no(s, *next_no())

# ============================================================ 9 两者关系：数据流
s = slide_base()
header(s, "RELATIONSHIP", "两系统的关系（一）：单向数据流", "数据只从互联网流向内网，永不开口回流 —— 这是安全模型的基石")
# 两大区块
rect(s, Inches(0.75), Inches(1.8), Inches(5.0), Inches(3.4), fill=CARD, line=RGBColor(0xE0, 0xC5, 0xC7), radius=0.06)
textbox(s, Inches(1.0), Inches(2.0), Inches(4.5), Inches(0.35),
        [("互联网监测平台", 16, CRIMSON, True, FONT_TITLE, 0)])
textbox(s, Inches(1.0), Inches(2.42), Inches(4.5), Inches(2.6),
        [("产出的数据资产：", 11, NAVY, True, FONT_BODY, 5),
         ("· 疫情事件（六源采集清洗后）", 10.5, GRAY, False, FONT_BODY, 4),
         ("· 国家风险评分与预警", 10.5, GRAY, False, FONT_BODY, 4),
         ("· 规则版本", 10.5, GRAY, False, FONT_BODY, 4),
         ("· 全球态势镜像数据", 10.5, GRAY, False, FONT_BODY, 8),
         ("摆渡前：打包 → gzip → SM4/AES 加密 → SM2/Ed25519 签名 → SHA-256 指纹", 10, CRIMSON, True, FONT_BODY, 0)])
# 箭头区
rect(s, Inches(6.0), Inches(2.85), Inches(1.35), Inches(1.3), fill=GOLD, radius=0.15)
textbox(s, Inches(6.05), Inches(3.1), Inches(1.25), Inches(0.9),
        [("单向", 15, PAPER, True, FONT_TITLE, 2),
         ("摆渡", 15, PAPER, True, FONT_TITLE, 2),
         ("ONE-WAY", 8, RGBColor(0xF3, 0xE6, 0xC8), False, FONT_LATIN, 0)], align=PP_ALIGN.CENTER)
textbox(s, Inches(5.95), Inches(2.5), Inches(1.5), Inches(0.3),
        [("▸▸▸", 18, GOLD, True, FONT_LATIN, 0)], align=PP_ALIGN.CENTER)
rect(s, Inches(7.6), Inches(1.8), Inches(5.0), Inches(3.4), fill=CARD, line=RGBColor(0xBE, 0xD9, 0xD4), radius=0.06)
textbox(s, Inches(7.85), Inches(2.0), Inches(4.5), Inches(0.35),
        [("内网口岸预警平台", 16, TEAL, True, FONT_TITLE, 0)])
textbox(s, Inches(7.85), Inches(2.42), Inches(4.5), Inches(2.6),
        [("接收后的加工：", 11, NAVY, True, FONT_BODY, 5),
         ("· 验签 → 解密 → Schema 校验 → 幂等入库", 10.5, GRAY, False, FONT_BODY, 4),
         ("· 离线镜像展示全球态势", 10.5, GRAY, False, FONT_BODY, 4),
         ("· 与本地旅客数据交叉匹配", 10.5, GRAY, False, FONT_BODY, 4),
         ("· 生成口岸布控建议", 10.5, GRAY, False, FONT_BODY, 8),
         ("接收通道：目录扫描 / RabbitMQ / 人工包导入 / API 轮询（默认关闭）", 10, TEAL, True, FONT_BODY, 0)])
# 底部原则
rect(s, Inches(0.75), Inches(5.45), Inches(11.85), Inches(1.3), fill=NAVY, radius=0.08)
textbox(s, Inches(1.05), Inches(5.65), Inches(11.3), Inches(0.95),
        [("单向性铁律", 12.5, PAPER, True, FONT_TITLE, 5),
         ("互联网端不会反向探测内网；内网默认零外联（无 API 轮询、无模型外呼）。跨平台 JWT 因 audience 隔离无法复用；数据库、Redis、密钥、构建产物全部独立 —— 即使一端失陷，另一端仍完整封闭。", 10.5, RGBColor(0x9F, 0xB4, 0xC2), False, FONT_BODY, 0)])
page_no(s, *next_no())

# ============================================================ 10 两者关系：隔离对照
s = slide_base()
header(s, "RELATIONSHIP", "两系统的关系（二）：全维度逻辑隔离", "同源同构的镜像体验，物理与逻辑的双重隔离")
rows = [
    ("维度", "互联网端", "内网端"),
    ("使命", "看世界：全球疫情监测与风险研判", "守国门：旅客风险匹配与口岸布控"),
    ("API 入口", "app.main:app / 8000", "app.intranet_main:app / 8002"),
    ("专属路由", "数据源、采集、风险重算、发送任务、outbox", "旅客、口岸建议、接收器、数据库直连"),
    ("前端入口", "internet.html / dist-internet", "intranet.html / dist-intranet"),
    ("数据库", "global_health_internet（独立 SQLite/PG）", "global_health_intranet（独立 SQLite/PG）"),
    ("缓存", "internet-redis", "intranet-redis"),
    ("容器网络", "internet_zone", "intranet_zone（internal，无外联）"),
    ("LLM 配置", "独立 llm_providers 表 + 字段加密", "独立 llm_providers 表 + 字段加密"),
    ("JWT", "audience=internet，禁止内网使用", "audience=intranet，禁止互联网使用"),
]
ty = Inches(1.8)
col_x = [Inches(0.75), Inches(2.6), Inches(7.3)]
col_w = [Inches(1.85), Inches(4.7), Inches(5.3)]
for ri, row in enumerate(rows):
    is_head = ri == 0
    rh = Inches(0.52) if is_head else Inches(0.5)
    for ci, cell in enumerate(row):
        fill = NAVY if is_head else (CARD if ri % 2 == 1 else LIGHT)
        fg = PAPER if is_head else (NAVY if ci == 0 else GRAY)
        rect(s, col_x[ci], ty, col_w[ci], rh, fill=fill)
        size = 11 if is_head else (10.5 if ci == 0 else 9.5)
        bold = True if (is_head or ci == 0) else False
        font = FONT_TITLE if is_head else (FONT_BODY if ci else FONT_BODY)
        textbox(s, col_x[ci] + Inches(0.15), ty + Inches(0.1), col_w[ci] - Inches(0.3), rh,
                [(cell, size, fg, bold, font, 0)])
    ty += rh + Inches(0.02)
textbox(s, Inches(0.75), ty + Inches(0.12), Inches(11.85), Inches(0.3),
        [("两平台共享的只有：设计语言、风险模型理念、规则树结构与摆渡数据 Schema —— 一切凭据与状态皆不共享。", 11, CRIMSON, True, FONT_BODY, 0)])
page_no(s, *next_no())

# ============================================================ 11 角色与权限
s = slide_base()
header(s, "ROLES & RBAC", "角色体系与最小权限", "五角色分工，职责分离，覆盖两平台全部操作面")
roles = [
    ("系统管理员", CRIMSON, ["用户与角色管理", "全量备份与恢复", "平台统计", "MFA 管理", "无审计日志查看权（职责分离）"]),
    ("数据分析员", TEAL, ["信息源配置与采集", "风险模型权重编辑", "规则编辑与发布", "摆渡任务管理"]),
    ("口岸操作员", GOLD, ["旅客录入与批量导入", "数据库直连同步", "风险匹配查看", "布控建议生成"]),
    ("审计员", NAVY, ["只读审计日志查询", "日志导出", "无任何写操作权限"]),
    ("只读用户", GRAY, ["全球态势查看", "事件与风险浏览", "地图图层查看"]),
]
cx = Inches(0.75)
cw = Inches(2.28); gap = Inches(0.13)
for t, c, items in roles:
    rect(s, cx, Inches(1.85), cw, Inches(3.6), fill=CARD, line=MIST, radius=0.08)
    rect(s, cx, Inches(1.85), cw, Inches(0.55), fill=c)
    textbox(s, cx, Inches(1.98), cw, Inches(0.35),
            [(t, 12.5, PAPER, True, FONT_BODY, 0)], align=PP_ALIGN.CENTER)
    lines = [("· " + it, 9.5, GRAY, False, FONT_BODY, 5) for it in items]
    textbox(s, cx + Inches(0.18), Inches(2.58), cw - Inches(0.32), Inches(2.8), lines)
    cx += cw + gap
rect(s, Inches(0.75), Inches(5.75), Inches(11.85), Inches(1.0), fill=LIGHT, radius=0.08)
textbox(s, Inches(1.05), Inches(5.95), Inches(11.3), Inches(0.65),
        [("账号安全基线", 12, NAVY, True, FONT_TITLE, 4),
         ("密码至少 10 位含大小写字母、数字与特殊字符；连续 5 次失败锁定 15 分钟；浏览器 30 分钟无操作会话超时；管理员首次登录须启用 TOTP 双因素。", 10, GRAY, False, FONT_BODY, 0)])
page_no(s, *next_no())

# ============================================================ 12 技术架构
s = slide_base()
header(s, "TECHNOLOGY", "技术架构与工程实践", "现代全栈 + 严格隔离 + 可验证")
tech = [
    ("后端", ["FastAPI (ASGI) 双入口", "SQLAlchemy 2.0 + Alembic 迁移", "PostgreSQL WAL / SQLite 开发库", "APScheduler 调度热更新", "pytest 覆盖率门禁 ≥70%"], CRIMSON),
    ("前端", ["React 18 + TypeScript", "Vite 双入口双构建产物", "MapLibre GL + Deck.gl 可视化", "本地 world-atlas 边界，零 CDN", "lucide 图标体系"], TEAL),
    ("数据与消息", ["RabbitMQ 摆渡消息通道", "Redis 缓存 + 故障降级", "gzip + SHA-256 指纹去重", "分片续传与幂等入库", "JSONL 结构化日志"], GOLD),
    ("部署", ["Docker Compose 双平台隔离部署", "离线镜像交付包（内网专用）", "网闸目录 / 单向 MQ / 白名单网关", "环境变量注入密钥", "WAF / 反向代理 / 堡垒机接入"], NAVY),
]
cx = Inches(0.75); cy = Inches(1.8); cw = Inches(2.88); chh = Inches(4.0)
for i, (t, items, c) in enumerate(tech):
    x = cx + i * (cw + Inches(0.13))
    rect(s, x, cy, cw, chh, fill=CARD, line=MIST, radius=0.07)
    rect(s, x, cy, cw, Inches(0.07), fill=c)
    textbox(s, x + Inches(0.2), cy + Inches(0.2), cw - Inches(0.4), Inches(0.35),
            [(t, 14, INK, True, FONT_TITLE, 0)])
    lines = [("· " + it, 9.8, GRAY, False, FONT_BODY, 6) for it in items]
    textbox(s, x + Inches(0.2), cy + Inches(0.62), cw - Inches(0.38), chh - Inches(0.8), lines)
rect(s, Inches(0.75), Inches(6.0), Inches(11.85), Inches(0.85), fill=NAVY, radius=0.08)
textbox(s, Inches(1.05), Inches(6.18), Inches(11.3), Inches(0.5),
        [("工程信条：互联网与内网使用不同 ASGI 入口、数据库、Redis、前端入口、构建产物和 Docker 镜像 —— 隔离即架构，镜像即体验。", 11.5, PAPER, True, FONT_BODY, 0)])
page_no(s, *next_no())

# ============================================================ 13 价值总结
s = slide_base()
header(s, "VALUE", "系统价值与建设成效", "从“被动看新闻”到“主动拦风险”的公共卫生治理升级")
vals = [
    ("早发现", "六源自动采集 + 趋势突变检测，全球疫情信号分钟级进入视野", CRIMSON),
    ("会研判", "六因子可解释评分 + 规则树版本管理，结论可追溯、可复核", TEAL),
    ("能落地", "风险结论单向摆渡至口岸一线，直接转化为布控动作", GOLD),
    ("守得住", "旅客数据字段加密全程脱敏；等保三级基线；国密可选", NAVY),
]
cx = Inches(0.75); cw = Inches(2.88)
for i, (t, d, c) in enumerate(vals):
    x = cx + i * (cw + Inches(0.13))
    rect(s, x, Inches(1.85), cw, Inches(2.3), fill=CARD, line=MIST, radius=0.08)
    rect(s, x, Inches(1.85), cw, Inches(0.65), fill=c)
    textbox(s, x, Inches(2.0), cw, Inches(0.4),
            [(t, 15, PAPER, True, FONT_TITLE, 0)], align=PP_ALIGN.CENTER)
    textbox(s, x + Inches(0.22), Inches(2.75), cw - Inches(0.44), Inches(1.3),
            [(d, 10.5, GRAY, False, FONT_BODY, 0)])
# 里程碑条
rect(s, Inches(0.75), Inches(4.45), Inches(11.85), Inches(2.2), fill=CARD, line=MIST, radius=0.06)
textbox(s, Inches(1.05), Inches(4.65), Inches(6.0), Inches(0.35),
        [("建设里程碑", 14, INK, True, FONT_TITLE, 0)])
mile = [("需求蓝图", "《001-需求》确立双平台边界"), ("平台落地", "采集/风险/地图/摆渡/旅客全功能实现"),
        ("智能增强", "大模型研判 + 连续测试 + 获取模型"), ("数据扩展", "MaxCompute / ADS 等广泛数据库直连")]
my2 = Inches(5.15)
for i, (t, d) in enumerate(mile):
    x = Inches(1.05) + i * Inches(2.93)
    rect(s, x, my2 + Inches(0.05), Inches(0.28), Inches(0.28), fill=CRIMSON if i % 2 == 0 else TEAL, radius=0.5)
    textbox(s, x + Inches(0.4), my2, Inches(2.4), Inches(0.3),
            [(t, 11.5, NAVY, True, FONT_BODY, 0)])
    textbox(s, x + Inches(0.4), my2 + Inches(0.32), Inches(2.5), Inches(0.6),
            [(d, 9.5, GRAY, False, FONT_BODY, 0)])
    if i < 3:
        rect(s, x + Inches(2.62), my2 + Inches(0.16), Inches(0.5), Inches(0.035), fill=MIST)
page_no(s, *next_no())

# ============================================================ 14 结尾
s = slide_base(NAVY)
rect(s, 0, Inches(7.28), SW, Inches(0.22), fill=CRIMSON)
textbox(s, Inches(1.2), Inches(2.3), Inches(10.9), Inches(1.4),
        [("互联网采全球之信号，内网守国门之一线", 30, PAPER, True, FONT_TITLE, 0)], align=PP_ALIGN.CENTER)
rect(s, SW/2 - Inches(0.75), Inches(3.85), Inches(1.5), Inches(0.045), fill=GOLD)
textbox(s, Inches(1.2), Inches(4.2), Inches(10.9), Inches(0.8),
        [("一套理念 · 两朵云 · 全链路闭环", 15, RGBColor(0x9F, 0xB4, 0xC2), True, FONT_BODY, 8),
         ("采集 — 研判 — 摆渡 — 匹配 — 布控，每一步皆可审计、可追溯、可验证", 12, RGBColor(0x7E, 0x93, 0xA3), False, FONT_BODY, 0)], align=PP_ALIGN.CENTER)
chip(s, SW/2 - Inches(2.6), Inches(5.6), Inches(1.7), Inches(0.45), "等保三级基线", RGBColor(0x3A, 0x53, 0x66), PAPER, 10.5)
chip(s, SW/2 - Inches(0.75), Inches(5.6), Inches(1.7), Inches(0.45), "国密可选套件", RGBColor(0x3A, 0x53, 0x66), PAPER, 10.5)
chip(s, SW/2 + Inches(1.1), Inches(5.6), Inches(1.7), Inches(0.45), "零外联内网", RGBColor(0x3A, 0x53, 0x66), PAPER, 10.5)
textbox(s, Inches(1.2), Inches(6.55), Inches(10.9), Inches(0.35),
        [("THANK YOU · 全球公共卫生监测与口岸预警系统", 11, RGBColor(0x7E, 0x93, 0xA3), False, FONT_LATIN, 0)], align=PP_ALIGN.CENTER)

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "全球公共卫生监测与口岸预警系统-双平台全景介绍.pptx")
prs.save(out)
print("saved:", out)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
